import os
from pypdf import PdfReader     # ← new import
from docx import Document as DocxDocument
from pptx import Presentation
import pdfplumber
import pytesseract
from PIL import Image
import re


class ContentExtractionService:

    @staticmethod
    def extract_text(file_path, file_type):
        file_type = file_type.lower()

        if file_type == 'pdf':
            text = ContentExtractionService._extract_pdf(file_path)

        elif file_type == 'docx':
            text = ContentExtractionService._extract_docx(file_path)

        elif file_type == 'pptx':
            text = ContentExtractionService._extract_pptx(file_path)

        elif file_type == 'txt':
            text = ContentExtractionService._extract_txt(file_path)

        else:
            raise Exception("Unsupported file type")

        # 🔥 APPLY CLEANING HERE
        text = ContentExtractionService._clean_text(text)

        return text
    # ---------------- PDF ----------------
    @staticmethod
    def _extract_pdf(file_path):
        text = ""

        try:
            # Try pdfplumber first
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    extracted = page.extract_text()
                    if extracted:
                        text += extracted + "\n"

            # If result is garbage (too many replacement chars) or too short → OCR
            garbage_ratio = text.count('\ufffd') / max(len(text), 1)
            if len(text.strip()) < 50 or garbage_ratio > 0.1:
                print("⚠️ Using OCR fallback")
                text = ""
                with pdfplumber.open(file_path) as pdf:
                    for page in pdf.pages:
                        img = page.to_image(resolution=300).original
                        ocr_text = pytesseract.image_to_string(img)
                        text += ocr_text + "\n"

        except Exception as e:
            print("❌ Extraction failed:", str(e))

        return text


    # ---------------- DOCX ----------------
    @staticmethod
    def _extract_docx(file_path):
        doc = DocxDocument(file_path)
        text = "\n".join([para.text for para in doc.paragraphs])
        return text

    # ---------------- PPTX ----------------
    @staticmethod
    def _extract_pptx(file_path):
        prs = Presentation(file_path)
        text = ""

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"

        return text

    # ---------------- TXT ----------------
    @staticmethod
    def _extract_txt(file_path):
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    
    def _clean_text(text):
        # Fix encoding issues but KEEP non-ASCII characters (don't strip them)
        text = text.encode("utf-8", "ignore").decode("utf-8")
        # Only remove actual control characters, not foreign scripts or symbols
        text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]', '', text)
        # Collapse excessive whitespace
        text = re.sub(r'\n{3,}', '\n\n', text)
        text = re.sub(r'[ \t]+', ' ', text)
        return text.strip()